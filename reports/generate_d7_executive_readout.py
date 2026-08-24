"""Generate the D7 executive readout from computed D2-D4 artifacts only."""

from __future__ import annotations

import json
import copy
import zipfile
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
REPORTS = ROOT / "reports"
PROCESSED = ROOT / "data" / "processed"
OUTPUT = REPORTS / "d7_executive_readout.pptx"

NAVY = RGBColor(17, 35, 54)
INK = RGBColor(30, 44, 58)
TEAL = RGBColor(16, 119, 124)
ORANGE = RGBColor(224, 116, 50)
PALE = RGBColor(239, 246, 245)
MUTED = RGBColor(92, 108, 119)
WHITE = RGBColor(255, 255, 255)

DISCLAIMER = "SYNTHETIC DEVELOPMENT DATA | Not official Zidio or NorthBay Living results"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def money(value: float | None) -> str:
    return "Not available" if value is None else f"₹{value:,.2f}"


def pct(value: float | None) -> str:
    return "Not available" if value is None else f"{value:.2%}"


def add_box(slide, left, top, width, height, fill=WHITE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, size=18, color=INK,
             bold=False, align=PP_ALIGN.LEFT, font="Aptos", margin=0.08):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, left, top, width, height, size=16, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
        p.bullet = True
    return box


def add_header(slide, title, subtitle=""):
    add_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.18), TEAL)
    add_text(slide, title, Inches(0.55), Inches(0.38), Inches(12.2), Inches(0.55), 28, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.58), Inches(0.95), Inches(12), Inches(0.35), 12, MUTED)
    add_text(slide, DISCLAIMER, Inches(0.58), Inches(7.12), Inches(12.1), Inches(0.22), 8, ORANGE, True)


def add_metric(slide, label, value, left, top, width=2.85, fill=PALE, value_color=NAVY):
    add_box(slide, left, top, width, Inches(1.2), fill, True)
    add_text(slide, label.upper(), left + Inches(0.12), top + Inches(0.12), width - Inches(0.24), Inches(0.25), 9, MUTED, True)
    add_text(slide, value, left + Inches(0.12), top + Inches(0.42), width - Inches(0.24), Inches(0.58), 22, value_color, True)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    return slide


def _normalise_pptx_metadata(path: Path) -> None:
    fixed_time = (2020, 1, 1, 0, 0, 0)
    temporary = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(temporary, "w") as target:
        for info in source.infolist():
            normalised = copy.copy(info)
            normalised.date_time = fixed_time
            target.writestr(normalised, source.read(info.filename))
    temporary.replace(path)


def main():
    d2 = load_json(REPORTS / "d2_eda_report.json")
    d3 = load_json(REPORTS / "d3_forecast_report.json")
    d4 = load_json(REPORTS / "d4_risk_report.json")
    recs = pd.read_csv(PROCESSED / "recommendations.csv")
    insights = d2["executive_summary"]["insights"]
    d4_counts = d4["decision_summary"]["counts"]
    rupee = d4["rupee_value_summary"]
    wape = d3["wape"]
    bias = d3["bias"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Executive summary
    slide = new_slide(prs)
    add_header(slide, "Executive Summary | FORESIGHT", "Inventory decisions from demand, forecast, and risk signals")
    add_text(slide, "A practical view of where inventory capital is exposed and what to review next.", Inches(0.65), Inches(1.55), Inches(7.3), Inches(0.7), 24, NAVY, True)
    add_box(slide, Inches(8.55), Inches(1.42), Inches(3.95), Inches(2.25), NAVY, True)
    add_text(slide, "VALUE AT STAKE", Inches(8.82), Inches(1.72), Inches(3.4), Inches(0.3), 11, RGBColor(178, 220, 218), True)
    add_text(slide, money(rupee["total_rupee_value_at_stake"]), Inches(8.78), Inches(2.05), Inches(3.5), Inches(0.75), 30, WHITE, True)
    add_text(slide, "Cost-basis exposure from D4", Inches(8.82), Inches(2.95), Inches(3.3), Inches(0.35), 12, RGBColor(220, 231, 235))
    add_metric(slide, "Reorder now", str(d4_counts.get("REORDER_NOW", 0)), Inches(0.68), Inches(3.0), 2.8, RGBColor(255, 241, 233), ORANGE)
    add_metric(slide, "Markdown / clear", str(d4_counts.get("MARKDOWN_CLEAR", 0)), Inches(3.75), Inches(3.0), 2.8, RGBColor(255, 246, 235), ORANGE)
    add_metric(slide, "Healthy", str(d4_counts.get("HEALTHY", 0)), Inches(6.82), Inches(3.0), 2.8, PALE, TEAL)
    add_bullets(slide, [
        "Protect availability for the 17 highest stockout-pressure SKUs.",
        "Review markdown or clearance for 7 excess-cover SKUs.",
        "Treat every figure as synthetic development evidence, not a company result.",
    ], Inches(0.7), Inches(4.65), Inches(11.8), Inches(1.45), 16)

    # 2. What is at risk
    slide = new_slide(prs)
    add_header(slide, "What Is At Risk?", "D4 separates availability pressure from excess inventory")
    add_metric(slide, "Total exposure", money(rupee["total_rupee_value_at_stake"]), Inches(0.7), Inches(1.45), 3.0, RGBColor(255, 241, 233), ORANGE)
    add_metric(slide, "Inventory exposure", money(rupee["total_inventory_value_exposure_rupees"]), Inches(3.95), Inches(1.45), 3.0)
    add_metric(slide, "Stockout shortfall", money(rupee["total_stockout_shortfall_rupees"]), Inches(7.2), Inches(1.45), 3.0)
    add_text(slide, "Decision mix", Inches(0.72), Inches(3.05), Inches(3), Inches(0.35), 18, NAVY, True)
    grid_items = [("REORDER_NOW", ORANGE), ("MARKDOWN_CLEAR", RGBColor(213, 151, 51)), ("WATCH_VOLATILE", RGBColor(180, 147, 52)), ("HEALTHY", TEAL)]
    for idx, (name, color) in enumerate(grid_items):
        y = 3.55 + idx * 0.65
        add_box(slide, Inches(0.75), Inches(y), Inches(0.16), Inches(0.34), color)
        add_text(slide, name.replace("_", " ").title(), Inches(1.05), Inches(y - 0.02), Inches(2.9), Inches(0.38), 14, INK, True)
        add_text(slide, str(d4_counts.get(name, 0)), Inches(3.65), Inches(y - 0.02), Inches(0.8), Inches(0.38), 16, NAVY, True, PP_ALIGN.RIGHT)
    add_bullets(slide, [
        f"{d4['stockout_risk_summary']['skus_with_shortfall']} SKUs show a computed lead-time shortfall.",
        f"{d4['overstock_risk_summary']['skus_with_excess']} SKUs exceed the healthy-cover inventory target.",
        "Rupee values use official unit_cost on a cost basis; lost revenue is not estimated.",
    ], Inches(5.1), Inches(3.35), Inches(7.2), Inches(2.4), 17)

    # 3. Recommended actions
    slide = new_slide(prs)
    add_header(slide, "Recommended Actions", "Prioritised directly from D4 output, highest value at stake first")
    add_text(slide, "Immediate review queue", Inches(0.72), Inches(1.38), Inches(4), Inches(0.35), 18, NAVY, True)
    top = recs.head(8)
    table = slide.shapes.add_table(len(top) + 1, 4, Inches(0.72), Inches(1.85), Inches(11.85), Inches(4.6)).table
    headers = ["SKU", "ACTION", "VALUE AT STAKE", "COVERAGE"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = WHITE
    for row_idx, (_, row) in enumerate(top.iterrows(), start=1):
        values = [row["sku_id"], str(row["decision"]).replace("_", " ").title(), money(row["rupee_value_at_stake"]), f"{row['coverage_weeks']:.1f} weeks"]
        for col, value in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = str(value)
            cell.fill.solid(); cell.fill.fore_color.rgb = PALE if row_idx % 2 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.color.rgb = INK

    # 4. Demand and inventory patterns
    slide = new_slide(prs)
    add_header(slide, "Demand & Inventory Patterns", "D2 observations connect demand concentration and seasonality to inventory review")
    add_box(slide, Inches(0.7), Inches(1.4), Inches(5.8), Inches(4.95), PALE, True)
    add_text(slide, "Computed observations", Inches(0.95), Inches(1.7), Inches(4.8), Inches(0.35), 18, NAVY, True)
    for idx, insight in enumerate(insights[:4]):
        add_text(slide, f"{idx + 1}. {insight['observation']}", Inches(0.95), Inches(2.2 + idx * 0.88), Inches(5.1), Inches(0.32), 13, INK, True)
        add_text(slide, insight["evidence"], Inches(1.12), Inches(2.52 + idx * 0.88), Inches(4.95), Inches(0.35), 11, MUTED)
    chart_paths = [REPORTS / "d2_charts" / "weekly_demand.png", REPORTS / "d2_charts" / "top_movers.png"]
    for idx, chart in enumerate(chart_paths):
        if chart.is_file():
            slide.shapes.add_picture(str(chart), Inches(6.85 + (idx % 2) * 2.95), Inches(1.65 + (idx // 2) * 2.5), width=Inches(2.7), height=Inches(2.15))
    add_text(slide, "Charts are descriptive summaries of synthetic development data.", Inches(6.9), Inches(6.05), Inches(5.3), Inches(0.3), 10, MUTED)

    # 5. Forecast performance
    slide = new_slide(prs)
    add_header(slide, "Forecast Performance", "D3 uses chronological rolling-origin validation")
    add_metric(slide, "Model WAPE", pct(wape["model"]), Inches(0.72), Inches(1.55), 2.8, PALE, TEAL)
    add_metric(slide, "Baseline WAPE", pct(wape["baseline"]), Inches(3.78), Inches(1.55), 2.8, RGBColor(241, 242, 244), NAVY)
    add_metric(slide, "WAPE improvement", pct(d3["model_vs_baseline"]["improvement_vs_baseline_wape"]), Inches(6.84), Inches(1.55), 2.8, RGBColor(232, 247, 239), TEAL)
    add_metric(slide, "Model bias", f"{bias['model']:.2f} units", Inches(9.9), Inches(1.55), 2.3, RGBColor(255, 246, 235), ORANGE)
    add_bullets(slide, [
        f"Selected model: {d3['selected_model'].replace('_', ' ').title()}.",
        f"Seasonal-naive baseline: {d3['seasonal_period_weeks']}-week lag.",
        f"{d3['backtest_methodology']['cv_folds']} rolling-origin folds with an {d3['forecast_horizon_weeks']}-week horizon.",
        "Negative bias means the forecasts under-predicted actual demand on average.",
        "WAPE is primary; bias is secondary. No random split was used.",
    ], Inches(0.85), Inches(3.35), Inches(11.3), Inches(2.4), 18)

    # 6. Operational meaning
    slide = new_slide(prs)
    add_header(slide, "What This Means Operationally", "Translate signals into review conversations, not automated actions")
    columns = [("Operations", ["Review REORDER_NOW SKUs first.", "Check supplier lead times and on-order quantities.", "Use coverage and shortfall together."]), ("Finance", ["Monitor cost-basis inventory exposure.", "Prioritise high-value markdown candidates.", "Separate cash tied up from service risk."]), ("Merchandising", ["Review seasonal and top-mover patterns.", "Investigate low-movement inventory.", "Treat volatility as a reason to review, not a forecast guarantee."])]
    for idx, (heading, bullets) in enumerate(columns):
        left = Inches(0.72 + idx * 4.15)
        add_box(slide, left, Inches(1.55), Inches(3.65), Inches(4.5), PALE if idx != 1 else RGBColor(255, 246, 235), True)
        add_text(slide, heading, left + Inches(0.22), Inches(1.87), Inches(3.15), Inches(0.45), 21, NAVY, True)
        add_bullets(slide, bullets, left + Inches(0.2), Inches(2.55), Inches(3.15), Inches(2.8), 15)

    # 7. Assumptions and limitations
    slide = new_slide(prs)
    add_header(slide, "Assumptions & Limitations", "What a decision-maker should know before interpreting the numbers")
    add_bullets(slide, [
        "All inputs are synthetic development data generated for testing; none are official internship or NorthBay Living extracts.",
        "D4 thresholds are implementation assumptions and must be reviewed with the mentor/client.",
        "Rupee exposure uses official unit_cost on a cost basis; lost revenue and margin impact are not inferred.",
        "Forecast performance is measured on this synthetic history and should be recomputed on official data.",
        "Promotion and seasonal relationships are descriptive, not causal claims.",
        "The system recommends review priorities; it does not place orders or take actions automatically.",
    ], Inches(0.95), Inches(1.55), Inches(11.3), Inches(4.7), 18)

    # 8. Next steps
    slide = new_slide(prs)
    add_header(slide, "Decision Summary & Next Steps", "A clean handoff from development evidence to official-data validation")
    add_text(slide, "Today", Inches(0.8), Inches(1.5), Inches(2), Inches(0.35), 18, TEAL, True)
    add_bullets(slide, ["Use the prioritised queue to demonstrate the workflow.", "Review the D3 accuracy and D4 valuation assumptions.", "Keep synthetic findings clearly separated from company conclusions."], Inches(0.85), Inches(1.95), Inches(5.2), Inches(2.2), 16)
    add_text(slide, "When official extracts arrive", Inches(7.0), Inches(1.5), Inches(4.5), Inches(0.35), 18, ORANGE, True)
    add_bullets(slide, ["Replace the four raw CSVs without changing the contract.", "Rerun D1 through D4 and compare data sufficiency, WAPE, bias, and risk exposure.", "Validate recommendations with Operations, Finance, and the mentor before use."], Inches(7.05), Inches(1.95), Inches(5.25), Inches(2.5), 16)
    add_box(slide, Inches(0.85), Inches(4.7), Inches(11.45), Inches(1.0), NAVY, True)
    add_text(slide, "The development pipeline is ready for an honest, reproducible replacement with official data.", Inches(1.15), Inches(4.95), Inches(10.85), Inches(0.5), 20, WHITE, True, PP_ALIGN.CENTER)

    prs.core_properties.created = datetime(2020, 1, 1)
    prs.core_properties.modified = datetime(2020, 1, 1)
    prs.save(OUTPUT)
    _normalise_pptx_metadata(OUTPUT)
    print(f"D7_OK slides={len(prs.slides)} output={OUTPUT}")


if __name__ == "__main__":
    main()
